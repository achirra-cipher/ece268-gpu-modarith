"""
Build 12_Project_Akhil_Sameera_Abhijit.pptx by adding slides into the class
template, following the template's own pattern:
  - Title slide keeps the course name as title; project title + names in subtitle
  - Content slides = title + a short body bullet list
All styling (fonts, colors, background) is inherited from the template as-is.

Usage: python3 make_pptx.py
"""
import sys
sys.path.insert(0, '/tmp/pptx_lib')

from pptx import Presentation

TEMPLATE = ("/Users/akhiltejachirra/Documents/UCSD/Quarter-3/ECE_268/Project/"
            "UCSD ECE 268 - Security of Hardware Embedded Systems.pptx")
OUTPUT   = ("/Users/akhiltejachirra/Documents/UCSD/Quarter-3/ECE_268/Project/"
            "slides/12_Project_Akhil_Sameera_Abhijit.pptx")

REL = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'

prs = Presentation(TEMPLATE)

# Map layouts by name (across both masters)
layouts = {}
for master in prs.slide_masters:
    for l in master.slide_layouts:
        layouts.setdefault(l.name, l)

# Remember the layouts the template's own slides used, then remove those slides
title_layout = None
body_layout = None
for sldId in list(prs.slides._sldIdLst):
    rId = sldId.get(REL)
    part = prs.part.related_part(rId)
    if part.slide_layout.name == 'TITLE':
        title_layout = part.slide_layout
    else:
        body_layout = part.slide_layout
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sldId)

title_layout = title_layout or layouts['TITLE']
body_layout  = body_layout  or layouts.get('TITLE_AND_BODY')


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def content_slide(title, bullets):
    slide = prs.slides.add_slide(body_layout)
    body = None
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = title
        elif body is None and idx != 0:
            body = ph
    if body is not None:
        tf = body.text_frame
        tf.clear()
        for i, line in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
    return slide


# ── SLIDE 1 — Title (template pattern: course name + subtitle) ───────────────
s = prs.slides.add_slide(title_layout)
for ph in s.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "GPU Modular Arithmetic with Barrett Reduction"
    elif ph.placeholder_format.idx == 1:
        ph.text = ("Group 12 — Akhil, Sameera, Abhijit\n"
                   "ECE 268 Final Project")
notes(s,
    "Hi everyone, we're Group 12 - Akhil, Sameera, and Abhijit. "
    "Our project implements a modular arithmetic engine on the GPU from scratch in CUDA, "
    "using Barrett reduction, and we apply it to the Number Theoretic Transform, "
    "which is a building block in post-quantum cryptography.")

# ── SLIDE 2 — Motivation ─────────────────────────────────────────────────────
s = content_slide("Motivation", [
    "Crypto like RSA, post-quantum NTT, and ZK-proofs run on millions of (a x b) mod p",
    "Doing \"mod p\" the normal way needs division, which is slow on GPUs",
    "Barrett reduction replaces the division with a multiply and a shift, using one precomputed constant",
    "The constant is the same for every thread, so it maps very well to the GPU",
])
notes(s,
    "Almost all of cryptography comes down to one operation repeated millions of times: "
    "a times b, mod p. The problem is that computing mod p normally needs an integer division, "
    "and division is slow on a GPU - it's not pipelined. "
    "Barrett reduction is a classic trick that replaces that division with just a multiply and a "
    "bit shift, using a constant we compute once ahead of time. "
    "Because that constant is shared across every thread, it fits the GPU's execution model really well. "
    "That's the core idea our whole project is built around.")

# ── SLIDE 3 — What we built ──────────────────────────────────────────────────
s = content_slide("What We Built", [
    "Modular multiply, exponentiation, and inverse, written from scratch in pure CUDA",
    "The same code runs on both GPU and CPU, which made testing straightforward",
    "Works over a 30-bit prime and a 64-bit prime field",
    "Used it to build a GPU Number Theoretic Transform (the core of lattice-based crypto)",
    "Checked every result against a CPU and a Python reference: no mismatches",
])
notes(s,
    "We wrote the modular arithmetic completely from scratch in CUDA - no libraries. "
    "One nice design choice: the same code compiles for both the GPU and the CPU, "
    "so we could test correctness on a laptop before running on the GPU. "
    "It supports modular multiplication, exponentiation, and inverse, over two prime fields - "
    "a 30-bit one and a 64-bit one. "
    "On top of that we built a GPU Number Theoretic Transform, which is the workhorse operation "
    "inside lattice-based post-quantum schemes like Kyber. "
    "We verified everything three ways - against a CPU version and an independent Python check - "
    "and got zero mismatches.")

# ── SLIDE 4 — Results ────────────────────────────────────────────────────────
s = content_slide("Results (Tesla T4 vs. CPU)", [
    "Modular multiplication: about 569x faster on the GPU",
    "Batched exponentiation, 1 million operations: 277x faster, all results correct",
    "NTT: up to 52x faster at large sizes",
    "At small sizes the GPU is slower - there isn't enough work to keep it busy",
])
notes(s,
    "Here are our results on a Tesla T4 GPU versus a single CPU core, running the exact same code. "
    "For raw modular multiplication, the GPU is about 569 times faster, because every multiply "
    "is independent and all the cores stay busy. "
    "For batched exponentiation - a million independent operations, which is the shape of bulk RSA "
    "or key exchange - we get 277 times faster, with every result correct. "
    "The NTT is up to 52 times faster at large sizes. "
    "But there's an honest caveat: at small sizes the GPU is actually slower than the CPU, "
    "because there just isn't enough work to fill it and the launch overhead dominates. "
    "The GPU only pays off once the problem is big enough.")

# ── SLIDE 5 — Takeaways ──────────────────────────────────────────────────────
s = content_slide("Takeaways", [
    "GPUs win when there are many independent operations to run at once",
    "Bigger primes mean more security but lower speed",
    "Our exponentiation isn't constant-time, so a real system would need that to be secure",
    "Barrett was a good fit: simple, fast, and no extra conversions needed",
])
notes(s,
    "A few takeaways to close. "
    "First, the GPU advantage really comes from having lots of independent work to do at once - "
    "that's where we saw the biggest speedups. "
    "Second, there's a clear trade-off: larger primes give more security but run slower, "
    "and that scales up to real sizes like RSA-2048. "
    "Third, and importantly for a security class: our exponentiation branches on the secret bits, "
    "so it's not constant-time - a real deployment would need to fix that to avoid timing attacks. "
    "And finally, Barrett reduction turned out to be a great fit because it's simple, fast, "
    "and doesn't need the extra conversions that the alternative, Montgomery, requires. "
    "Thanks - happy to take any questions.")

prs.save(OUTPUT)
print(f"Saved -> {OUTPUT}")
print(f"Title layout used: {title_layout.name}; body layout used: {body_layout.name}")
